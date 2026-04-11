"""
Format-preserving document translation endpoints.

Translates DOCX, XLSX, and PPTX files while preserving all formatting,
styles, fonts, colors, merged cells, formulas, slide layouts, etc.

Uses the same AI translation pipeline as tilshunos_routes.translate_text.
"""
import io
import os
import re
import logging
import tempfile
from typing import Dict, Any

from fastapi import APIRouter, UploadFile, File, Form, HTTPException, Depends
from fastapi.responses import StreamingResponse

from auth import get_current_user

logger = logging.getLogger("document_routes")
router = APIRouter(prefix="/api/documents", tags=["documents"])

LANG_LABELS = {
    "en": "English",
    "ru": "Russian",
    "uz": "Uzbek (Latin)",
    "uz-lat": "Uzbek (Latin)",
    "uz-cyr": "Uzbek (Cyrillic)",
}


async def _translate_text(text: str, source_lang: str, target_lang: str) -> str:
    """Translate a single text string using the AI pipeline (same as tilshunos)."""
    if not text or not text.strip():
        return text

    # Skip texts that are just numbers, punctuation, or whitespace
    cleaned = re.sub(r'[\d\s\W]+', '', text)
    if not cleaned:
        return text

    base_src = "ru" if source_lang == "ru" else ("uz" if source_lang.startswith("uz") else "en")
    base_tgt = "ru" if target_lang == "ru" else ("uz" if target_lang.startswith("uz") else "en")

    if base_src == base_tgt:
        return text

    # 1. TM check
    try:
        import tm_search
        tm_match = tm_search.best_match(text, src_lang=base_src, tgt_lang=base_tgt)
        if tm_match and tm_match["score"] >= 0.80:
            return tm_match["tgt_text"]
    except Exception:
        pass

    # 2. Tahrirchi dilmash
    try:
        import tahrirchi_engine
        if tahrirchi_engine.is_available():
            tr = await tahrirchi_engine.translate_async(text, source_lang, target_lang)
            if tr and len(tr.strip()) > 1 and tr.strip() != text.strip():
                return tr.strip()
    except Exception:
        pass

    # 3. NLLB for Russian pairs
    needs_russian = "ru" in (base_src, base_tgt)
    if needs_russian:
        try:
            import translator_engine
            if translator_engine.is_available() and translator_engine.supports(source_lang, target_lang):
                tr = await translator_engine.translate_async(text, source_lang, target_lang)
                if tr and len(tr.strip()) > 1:
                    return tr.strip()
        except Exception:
            pass

    # 4. Llama
    try:
        import llama_engine
        if llama_engine.is_available():
            tr = await llama_engine.translate(text, source_lang, target_lang)
            if tr and len(tr.strip()) > 1:
                return tr.strip()
    except Exception:
        pass

    # 5. Mistral
    try:
        import mistral_engine
        if mistral_engine.is_available():
            tr = await mistral_engine.translate(text, source_lang, target_lang)
            if tr and len(tr.strip()) > 1:
                return tr.strip()
    except Exception:
        pass

    # 6. Cloud fallback (Gemini / Claude)
    try:
        from routes.ai_helpers import generate_ai_content
        src_label = LANG_LABELS.get(source_lang, source_lang)
        tgt_label = LANG_LABELS.get(target_lang, target_lang)
        prompt = f"Translate the following {src_label} text to {tgt_label}. Return ONLY the translation, no explanations.\n\n{text}"
        translated = await generate_ai_content(prompt, prefer="cloud")
        result = (translated or "").strip()
        if result:
            # Save to TM for future reuse
            try:
                import tm_search
                tm_search.save_to_tm(text, result, src_lang=base_src, tgt_lang=base_tgt, source="doc_translate")
            except Exception:
                pass
            return result
    except Exception as e:
        logger.warning(f"[doc-translate] Cloud fallback failed: {e}")

    return text  # Return original if all engines fail


async def _translate_batch(texts: list[str], source_lang: str, target_lang: str) -> list[str]:
    """Translate a batch of texts. Processes sequentially to respect rate limits."""
    results = []
    for t in texts:
        translated = await _translate_text(t, source_lang, target_lang)
        results.append(translated)
    return results


# ═══════════════════════════════════════════════════
# DOCX Translation
# ═══════════════════════════════════════════════════

@router.post("/translate-docx")
async def translate_docx(
    file: UploadFile = File(...),
    source_lang: str = Form("en"),
    target_lang: str = Form("uz-lat"),
    engine: str = Form("auto"),
    current_user: Dict = Depends(get_current_user),
):
    """Translate a DOCX file preserving all formatting (fonts, bold, italic, sizes, colors)."""
    if not file.filename or not file.filename.lower().endswith('.docx'):
        raise HTTPException(status_code=400, detail="Only .docx files are accepted")

    try:
        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
    except ImportError:
        raise HTTPException(status_code=500, detail="python-docx is not installed")

    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Empty file")

    try:
        doc = Document(io.BytesIO(raw))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid DOCX file: {e}")

    # Translate paragraphs — iterate runs to preserve formatting
    for para in doc.paragraphs:
        for run in para.runs:
            if run.text and run.text.strip():
                run.text = await _translate_text(run.text, source_lang, target_lang)

    # Translate text in tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if run.text and run.text.strip():
                            run.text = await _translate_text(run.text, source_lang, target_lang)

    # Translate headers and footers
    for section in doc.sections:
        for header_footer in [section.header, section.footer]:
            if header_footer and header_footer.is_linked_to_previous is False:
                for para in header_footer.paragraphs:
                    for run in para.runs:
                        if run.text and run.text.strip():
                            run.text = await _translate_text(run.text, source_lang, target_lang)

    # Save to buffer
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    out_name = file.filename.rsplit('.', 1)[0] + f"_{target_lang}.docx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{out_name}"'},
    )


# ═══════════════════════════════════════════════════
# XLSX Translation
# ═══════════════════════════════════════════════════

@router.post("/translate-xlsx")
async def translate_xlsx(
    file: UploadFile = File(...),
    source_lang: str = Form("en"),
    target_lang: str = Form("uz-lat"),
    engine: str = Form("auto"),
    current_user: Dict = Depends(get_current_user),
):
    """Translate an XLSX file preserving formatting, formulas, merged cells, and styles."""
    if not file.filename or not file.filename.lower().endswith('.xlsx'):
        raise HTTPException(status_code=400, detail="Only .xlsx files are accepted")

    try:
        from openpyxl import load_workbook
    except ImportError:
        raise HTTPException(status_code=500, detail="openpyxl is not installed")

    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Empty file")

    try:
        wb = load_workbook(io.BytesIO(raw))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid XLSX file: {e}")

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                # Skip formulas (they start with =)
                if cell.value and isinstance(cell.value, str):
                    val = cell.value.strip()
                    if val and not val.startswith('='):
                        # Skip pure numbers disguised as strings
                        try:
                            float(val.replace(',', '').replace(' ', ''))
                            continue
                        except ValueError:
                            pass
                        cell.value = await _translate_text(cell.value, source_lang, target_lang)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)

    out_name = file.filename.rsplit('.', 1)[0] + f"_{target_lang}.xlsx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{out_name}"'},
    )


# ═══════════════════════════════════════════════════
# PPTX Translation
# ═══════════════════════════════════════════════════

@router.post("/translate-pptx")
async def translate_pptx(
    file: UploadFile = File(...),
    source_lang: str = Form("en"),
    target_lang: str = Form("uz-lat"),
    engine: str = Form("auto"),
    current_user: Dict = Depends(get_current_user),
):
    """Translate a PPTX file preserving slide layout, fonts, colors, and positioning."""
    if not file.filename or not file.filename.lower().endswith('.pptx'):
        raise HTTPException(status_code=400, detail="Only .pptx files are accepted")

    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx is not installed")

    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Empty file")

    try:
        prs = Presentation(io.BytesIO(raw))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid PPTX file: {e}")

    for slide in prs.slides:
        for shape in slide.shapes:
            # Text frames (titles, body text, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text and run.text.strip():
                            run.text = await _translate_text(run.text, source_lang, target_lang)

            # Tables inside slides
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if run.text and run.text.strip():
                                        run.text = await _translate_text(run.text, source_lang, target_lang)

        # Slide notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip():
                        run.text = await _translate_text(run.text, source_lang, target_lang)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    out_name = file.filename.rsplit('.', 1)[0] + f"_{target_lang}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{out_name}"'},
    )


# ═══════════════════════════════════════════════════
# Preview endpoint (extract text summary for UI preview)
# ═══════════════════════════════════════════════════

@router.post("/preview")
async def preview_document(
    file: UploadFile = File(...),
    current_user: Dict = Depends(get_current_user),
):
    """Extract text preview from a document (DOCX/XLSX/PPTX) for the UI."""
    name = (file.filename or "").lower()
    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Empty file")

    paragraphs = []
    doc_type = "unknown"

    if name.endswith('.docx'):
        doc_type = "docx"
        from docx import Document
        doc = Document(io.BytesIO(raw))
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text)

    elif name.endswith('.xlsx'):
        doc_type = "xlsx"
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(raw), read_only=True)
        for ws in wb.worksheets:
            for row in ws.iter_rows(max_row=100, values_only=True):
                for val in row:
                    if val and isinstance(val, str) and val.strip():
                        paragraphs.append(val)
        wb.close()

    elif name.endswith('.pptx'):
        doc_type = "pptx"
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            paragraphs.append(para.text)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type. Use .docx, .xlsx, or .pptx")

    return {
        "type": doc_type,
        "filename": file.filename,
        "paragraphs": paragraphs[:200],  # Limit preview to 200 paragraphs
        "total_paragraphs": len(paragraphs),
    }
